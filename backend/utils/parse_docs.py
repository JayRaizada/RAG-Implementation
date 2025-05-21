from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

def parse_document(file_path):
    if file_path.endswith(".pdf"):
        reader = PdfReader(file_path)
        return "\n".join([page.extract_text() for page in reader.pages])
    elif file_path.endswith(".docx"):
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    elif file_path.endswith(".pptx"):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        print(text)
        return text
    return ""
